"""
Hapi Fintech Case Study – Premium Presentation v2
Design inspired by: clean consultant decks (white bg, accent bars, dark section slides)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# ── Palette (matching reference design with Hapi brand) ─────────────────────
PUR     = RGBColor(0x6B, 0x3F, 0xFF)   # Hapi purple
TEAL    = RGBColor(0x00, 0xC2, 0xB5)   # Hapi teal
ORANGE  = RGBColor(0xF0, 0x71, 0x32)   # accent orange (from reference)
DARK    = RGBColor(0x35, 0x38, 0x4E)   # dark section bg (charcoal)
NEAR_BK = RGBColor(0x1A, 0x1A, 0x2E)   # near-black text
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF5, 0xF6, 0xFA)   # slide background off-white
MID_GY  = RGBColor(0xCC, 0xCC, 0xCC)
TEXT_GY = RGBColor(0x55, 0x55, 0x66)   # body text gray
GREEN   = RGBColor(0x28, 0xA7, 0x45)
RED     = RGBColor(0xDC, 0x35, 0x45)
YELLOW  = RGBColor(0xFF, 0xC1, 0x07)

BASE = "/home/user/clase13"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════════
# LOW-LEVEL HELPERS
# ══════════════════════════════════════════════════════════════════

def rgb(r, g, b):
    return RGBColor(r, g, b)

def rect(slide, l, t, w, h, fill=None, line=None, lw=0, radius=0):
    """Add a rectangle (optionally rounded via XML)."""
    shape = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h, size=14, bold=False, color=NEAR_BK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True, spacing=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.space_before = Pt(spacing)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_image(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))


def accent_bars(slide, l=0.38, t=1.05, bar_w=0.33, bar_h=0.055):
    """Two small accent bars (teal + orange) above a title — like reference style."""
    rect(slide, l, t, bar_w, bar_h, fill=TEAL)
    rect(slide, l + bar_w + 0.04, t, bar_w, bar_h, fill=ORANGE)


def slide_header(slide, label=""):
    """Minimal top strip: thin line + small label."""
    rect(slide, 0, 0, 13.33, 0.42, fill=WHITE)
    rect(slide, 0, 0.42, 13.33, 0.02, fill=MID_GY)
    if label:
        txt(slide, label.upper(), 0.38, 0.06, 8, 0.3,
            size=8, color=TEXT_GY, italic=True)
    txt(slide, "hapi · Case Study 2026", 11.5, 0.06, 1.8, 0.3,
        size=8, color=TEXT_GY, align=PP_ALIGN.RIGHT)


def white_slide(slide, title, section_label=""):
    """Standard content slide: white bg, accent bars, large title."""
    rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(slide, section_label)
    accent_bars(slide, l=0.38, t=1.1)
    txt(slide, title, 0.38, 1.22, 12.5, 0.95,
        size=32, bold=True, color=NEAR_BK)


def dark_section(slide, section_label, headline):
    """Dark charcoal section-divider slide."""
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    # Left purple stripe
    rect(slide, 0, 0, 0.08, 7.5, fill=PUR)
    accent_bars(slide, l=0.38, t=0.95, bar_w=0.28, bar_h=0.05)
    txt(slide, section_label, 0.38, 1.1, 10, 0.35,
        size=11, color=TEAL, bold=True)
    txt(slide, headline, 0.38, 1.55, 10.5, 2.5,
        size=46, bold=True, color=WHITE)


def circle_num(slide, num, l, t, r=0.28, color=TEAL):
    """Draw a filled circle with a number inside."""
    rect(slide, l, t, r*2, r*2, fill=color)
    txt(slide, str(num), l, t + 0.02, r*2, r*2 - 0.05,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def bullet_row(slide, num, text, l, t, w, color=TEAL):
    circle_num(slide, num, l, t + 0.03, color=color)
    txt(slide, text, l + 0.72, t, w - 0.72, 0.55,
        size=12, color=NEAR_BK)


def stat_box(slide, value, label, sub, l, t, w=2.5, h=1.45, fill=WHITE, vc=PUR):
    rect(slide, l, t, w, h, fill=fill, line=rgb(220,222,230), lw=0.8)
    txt(slide, value, l+0.1, t+0.1, w-0.2, 0.7,
        size=28, bold=True, color=vc, align=PP_ALIGN.CENTER)
    txt(slide, label, l+0.1, t+0.78, w-0.2, 0.35,
        size=11, bold=True, color=NEAR_BK, align=PP_ALIGN.CENTER)
    if sub:
        txt(slide, sub, l+0.1, t+1.1, w-0.2, 0.3,
            size=9, color=TEXT_GY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=rgb(13, 17, 40))

# Full-width bottom image (chart bg)
add_image(s, f"{BASE}/bg_title.png", 0, 3.1, 13.33, 4.4)

# Semi-transparent overlay on top half
rect(s, 0, 0, 13.33, 4.5, fill=rgb(13, 17, 40))

# Hapi logo top-left
add_image(s, f"{BASE}/hapi_logo.png", 0.45, 0.28, 3.2, 0.95)

# Top-right metadata
txt(s, "Final Case Study  ·  Fintech Seminar  ·  June 2026",
    8.5, 0.35, 4.6, 0.38, size=9, color=rgb(140,140,180),
    align=PP_ALIGN.RIGHT, italic=True)

# Purple left accent bar
rect(s, 0, 1.55, 0.07, 2.1, fill=PUR)

# Main headline
txt(s, "Democratizing\nInvesting in\nLatin America",
    0.55, 1.6, 8.5, 2.4, size=48, bold=True, color=WHITE)

# Tagline
txt(s, "Access the US market from just US$5 · 500K+ users · 20+ countries",
    0.55, 4.2, 9.5, 0.5, size=14, color=TEAL, italic=True)

# Divider
rect(s, 0.55, 4.85, 5.5, 0.035, fill=PUR)

# Team names
txt(s, "Daniel Taboada  ·  Benjamin Pereyra  ·  Luis Sanchez\nJorge Ramirez  ·  Santiago Wiesse",
    0.55, 5.0, 9, 0.7, size=12, color=rgb(180, 180, 210))
txt(s, "Universidad del Pacífico  ·  Fanny Mariana Díaz Galindo",
    0.55, 5.75, 9, 0.4, size=11, color=rgb(130, 130, 165), italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=TEAL)

# Left panel (darker teal)
rect(s, 0, 0, 5.2, 7.5, fill=rgb(0, 160, 150))
txt(s, "Agenda", 0.45, 0.9, 4.3, 1.0,
    size=42, bold=True, color=WHITE)
txt(s, "Hapi · Case Study 2026", 0.45, 2.0, 4.3, 0.4,
    size=11, color=rgb(200, 245, 242), italic=True)

# Logo on left panel
add_image(s, f"{BASE}/hapi_logo.png", 0.45, 6.5, 2.6, 0.78)

# Right: agenda items as a grid
items = [
    ("01", "Product & Technology"),
    ("02", "Market Opportunity"),
    ("03", "Competitive Landscape"),
    ("04", "Business Model & Revenue"),
    ("05", "Customer Segments"),
    ("06", "Key Metrics & Funding"),
    ("07", "SWOT Analysis"),
    ("08", "Strategic Outlook"),
]
cols = 2
iw = 3.8
ih = 0.72
gap_x = 0.18
gap_y = 0.22
sl = 5.5
st = 0.85

for i, (num, title) in enumerate(items):
    col = i % cols
    row = i // cols
    l = sl + col * (iw + gap_x)
    t = st + row * (ih + gap_y)
    rect(s, l, t, iw, ih, fill=rgb(0,170,160))
    txt(s, num, l + 0.15, t + 0.12, 0.5, 0.5,
        size=18, bold=True, color=rgb(200,245,242))
    txt(s, title, l + 0.65, t + 0.16, iw - 0.8, 0.45,
        size=13, bold=False, color=WHITE)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — DARK: Product & Technology intro
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_section(s, "01 · Product & Technology",
             "Built by investors,\nfor investors.")
# Right side image (mobile mockup)
add_image(s, f"{BASE}/bg_mobile.png", 8.0, 0.5, 5.0, 6.8)

txt(s, "Three ex-UP students couldn't invest in US stocks without\nhigh minimums, fees, and paperwork — so they built Hapi.",
    0.45, 4.3, 7.2, 0.9, size=13, color=rgb(200, 200, 220), italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — PRODUCT FEATURES
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_slide(s, "Platform at a Glance", "01 · Product & Technology")

# Stat boxes — top row
stats = [
    ("2020", "Founded", "By Dusko Kelez,\nPiero Sifuentes,\nBilly Caballero"),
    ("500K+", "Users", "across 20+ countries"),
    ("37", "Employees", "lean, outsourced\ninfra model"),
    ("12,000+", "Instruments", "stocks, ETFs & crypto"),
    ("US$5", "Min. Investment", "fractional shares\nincluded"),
]
sw = 2.38
sg = 0.14
sl4 = 0.35
for i, (v, l, s4) in enumerate(stats):
    stat_box(s, v, l, s4, sl4 + i*(sw+sg), 2.4, w=sw, h=1.5, vc=PUR)

# Feature bullets
rect(s, 0.35, 4.15, 12.6, 0.04, fill=TEAL)
txt(s, "Key Features", 0.35, 4.28, 5, 0.38,
    size=14, bold=True, color=NEAR_BK)

feats = [
    ("Stocks & ETFs", "Zero commissions on all trades, 12,000+ instruments"),
    ("Fractional shares", "Start with as little as US$5 — no high minimums"),
    ("Crypto", "From US$1 via Bakkt Crypto Solutions"),
    ("DRIP", "Auto-reinvest dividends — hands-free compounding"),
    ("Hapi Prime", "US$9.99/mo — real-time prices, instant settlement, alerts"),
]
fw = 4.0
for i, (feat, desc) in enumerate(feats):
    col = i % 3
    row = i // 3
    fl = 0.35 + col * (fw + 0.2)
    ft = 4.75 + row * 0.75
    rect(s, fl, ft, fw, 0.65, fill=LIGHT, line=rgb(220,222,235), lw=0.5)
    txt(s, feat, fl + 0.15, ft + 0.07, 1.5, 0.28, size=11, bold=True, color=PUR)
    txt(s, desc, fl + 1.7, ft + 0.07, fw - 1.9, 0.5, size=10, color=NEAR_BK)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — TECHNOLOGY & COMPLIANCE (split layout)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_slide(s, "Infrastructure & Compliance", "01 · Product & Technology")

# Left col: tech
txt(s, "Technology Stack", 0.35, 2.3, 6.0, 0.4, size=14, bold=True, color=NEAR_BK)
tech = [
    "Clearing & Custody → Apex Clearing Corp\n   (same infra as major US brokers)",
    "Crypto → Bakkt Crypto Solutions",
    "2-Factor Authentication + end-to-end encryption",
    "Only 37 employees — outsourcing enables scale",
]
for i, t5 in enumerate(tech):
    bullet_row(s, i+1, t5, 0.35, 2.82 + i * 0.95, 6.0, color=TEAL)

# Right col: compliance badges
rx = 7.0
rect(s, rx, 2.15, 5.95, 4.85, fill=LIGHT, line=rgb(220,222,235), lw=0.8)
txt(s, "Regulatory Protection", rx+0.2, 2.28, 5.5, 0.4,
    size=14, bold=True, color=NEAR_BK)

badges = [
    ("SEC Registered", "Hapi Securities LLC", PUR),
    ("FINRA Member",   "US broker-dealer rules", TEAL),
    ("SIPC Coverage",  "Up to US$500,000 per account", GREEN),
    ("⚠  Crypto",      "NOT covered by SIPC", RED),
]
for i, (b, d, c) in enumerate(badges):
    bt = 2.82 + i * 1.06
    rect(s, rx+0.2, bt, 1.9, 0.42, fill=c)
    txt(s, b, rx+0.28, bt+0.07, 1.75, 0.32,
        size=11, bold=True, color=WHITE)
    txt(s, d, rx+2.3, bt+0.1, 3.4, 0.32, size=11, color=NEAR_BK)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — DARK: Market Opportunity intro
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_section(s, "02 · Market Opportunity",
             "A $15B market\nstill untapped.")
# LATAM bg on right
add_image(s, f"{BASE}/bg_latam.png", 7.5, 0.8, 5.5, 6.2)

txt(s, "~90% of banking assets in 5 institutions.\nSmall investors ignored. Local currencies losing value.\nHapi fixes all three.",
    0.45, 4.5, 6.8, 1.2, size=13, color=rgb(200,200,220), italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — MARKET OPPORTUNITY (data)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_slide(s, "Market Opportunity", "02 · Market Opportunity")

# 3 big stat cards
big_stats = [
    ("US$15.23B", "LATAM Fintech Market", "IMARC Group", PUR),
    ("~90%", "Banking controlled by 5 institutions", "Mordor Intelligence", RED),
    ("1M", "Users targeted by Hapi", "CEO target, Gestión 2025", TEAL),
]
for i, (v, l, src, c) in enumerate(big_stats):
    bl = 0.35 + i * 4.3
    rect(s, bl, 2.3, 4.05, 1.8, fill=c)
    txt(s, v, bl+0.15, 2.38, 3.75, 0.9,
        size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, l, bl+0.15, 3.28, 3.75, 0.5,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, src, bl+0.15, 3.8, 3.75, 0.25,
        size=9, color=rgb(200,230,230) if c==TEAL else WHITE,
        italic=True, align=PP_ALIGN.CENTER)

# Tailwinds
rect(s, 0.35, 4.35, 12.6, 0.04, fill=TEAL)
txt(s, "Why the timing is right", 0.35, 4.48, 5, 0.35, size=13, bold=True, color=NEAR_BK)

tailwinds = [
    "Rising smartphone & internet penetration across LATAM",
    "Local currencies losing value → demand for USD-denominated assets",
    "Growing financial literacy — new generation of retail investors",
    "Traditional banks never cared about small investors — structural gap",
]
for i, tw in enumerate(tailwinds):
    col = i % 2
    row = i // 2
    tl = 0.35 + col * 6.4
    tt = 4.93 + row * 0.65
    rect(s, tl, tt, 6.2, 0.55, fill=LIGHT, line=rgb(220,222,235), lw=0.5)
    txt(s, "→  " + tw, tl+0.15, tt+0.1, 5.9, 0.42, size=11, color=NEAR_BK)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — COMPETITIVE LANDSCAPE
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_slide(s, "Competitive Landscape", "03 · Competitive Landscape")

# Table
headers = ["Platform", "Min.", "Commission", "SIPC", "LATAM\nFocus", "Fractional", "Education"]
rows = [
    ["Hapi",                "US$5",  "Zero",  "✔ $500K", "✔ Primary",  "✔",  "✔"],
    ["Interactive Brokers", "US$0",  "Low",   "✔ $500K", "Partial",    "✔",  "✘"],
    ["Charles Schwab",      "US$0",  "Zero",  "✔ $500K", "✘",          "✔",  "Partial"],
    ["XTB",                 "US$0",  "Zero*", "✘",       "Partial",    "✔",  "✔"],
    ["eToro",               "US$10", "Zero*", "✘",       "Partial",    "✔",  "✔"],
    ["Trii / Tyba",         "Low",   "Low",   "✘",       "✔ Local",    "Ltd","✔"],
]

col_w = [2.4, 1.0, 1.35, 1.25, 1.35, 1.2, 1.25]
col_l = [0.35]
for w in col_w[:-1]:
    col_l.append(col_l[-1] + w + 0.04)

row_h = 0.52
ht = 2.28
dt = ht + row_h + 0.04

# Header row
for ci, (h, cw, cl) in enumerate(zip(headers, col_w, col_l)):
    rect(s, cl, ht, cw, row_h, fill=NEAR_BK)
    txt(s, h, cl+0.05, ht+0.05, cw-0.1, row_h-0.1,
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Data rows
for ri, row in enumerate(rows):
    rt = dt + ri * (row_h + 0.04)
    is_hapi = row[0] == "Hapi"
    fill_c = rgb(235, 230, 255) if is_hapi else WHITE
    lc = PUR if is_hapi else rgb(220,222,235)
    lw = 1.5 if is_hapi else 0.5
    for ci, (cell, cw, cl) in enumerate(zip(row, col_w, col_l)):
        rect(s, cl, rt, cw, row_h, fill=fill_c, line=lc, lw=lw)
        tc = PUR if is_hapi else (GREEN if "✔" in cell else
                                   RED if "✘" in cell else NEAR_BK)
        txt(s, cell, cl+0.05, rt+0.1, cw-0.1, row_h-0.12,
            size=10, bold=is_hapi, color=tc, align=PP_ALIGN.CENTER)

txt(s, "* Zero commission on stocks, but spreads may apply",
    0.35, 6.85, 8, 0.35, size=9, color=TEXT_GY, italic=True)

# Competitive chart
add_image(s, f"{BASE}/chart_competitive.png", 9.0, 2.25, 4.2, 4.5)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — DARK: Business Model intro
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_section(s, "04 · Business Model",
             "Zero commissions.\nMultiple revenue streams.")
txt(s, "\"The company makes money in several ways —\nso it doesn't depend on just one.\"",
    0.45, 4.5, 7.5, 1.1, size=14, color=TEAL, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — REVENUE STREAMS
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_slide(s, "Revenue Streams", "04 · Business Model")

streams = [
    ("Payment for\nOrder Flow", "Fees from routing orders\nto exchanges (PFOF)", PUR, "~35%"),
    ("Hapi Prime\nUS$9.99/mo", "Real-time prices, instant\nsettlement, alerts", TEAL, "~25%"),
    ("Deposit &\nWithdrawal", "0.9–1% on deposits\nUS$4.99–$10 withdrawals", ORANGE, "~15%"),
    ("Idle Cash\nInterest", "Yield on uninvested\ncash in accounts", GREEN, "~12%"),
    ("Crypto\nSpread", "1% via Bakkt on\nall crypto trades", rgb(100,60,200), "~8%"),
    ("Closing\nFees", "US$0.10 – $0.15\nper trade close", rgb(100,100,130), "~5%"),
]

sw6 = 2.04
sg6 = 0.08
sl6 = 0.35
st6 = 2.3
sh6 = 2.25

for i, (title, desc, c, pct) in enumerate(streams):
    l = sl6 + i * (sw6 + sg6)
    # Color header
    rect(s, l, st6, sw6, 0.52, fill=c)
    txt(s, title, l+0.08, st6+0.06, sw6-0.16, 0.44,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # White body
    rect(s, l, st6+0.52, sw6, sh6-0.52, fill=WHITE, line=c, lw=1.2)
    txt(s, desc, l+0.1, st6+0.62, sw6-0.2, 0.7,
        size=11, color=NEAR_BK, align=PP_ALIGN.CENTER)
    # Percentage badge
    rect(s, l+sw6/2-0.38, st6+1.38, 0.76, 0.44, fill=c)
    txt(s, pct, l+sw6/2-0.38, st6+1.44, 0.76, 0.36,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Revenue chart
add_image(s, f"{BASE}/chart_revenue.png", 4.0, 4.65, 5.5, 2.7)

# Key insight
rect(s, 0.35, 4.7, 3.5, 1.9, fill=LIGHT, line=PUR, lw=1.2)
txt(s, "Key Model Insight", 0.5, 4.83, 3.2, 0.38, size=12, bold=True, color=PUR)
txt(s, "Average position:\n~US$250 per user\n\nMonetizes through\nvolume & scale —\nnot trade size",
    0.5, 5.25, 3.2, 1.28, size=11, color=NEAR_BK)

txt(s, "⚠  PFOF controversy: SEC tightened disclosure rules (605/606). Restricted in EU by ESMA.",
    0.35, 6.85, 12.6, 0.38, size=9, color=TEXT_GY, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — CUSTOMER SEGMENTS (image + list layout)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_slide(s, "Customer Segments", "05 · Customer Segments")

segments = [
    ("01", "Beginner\nInvestors",
     ["Everyday retail users,\nstudents, young professionals",
      "Zero-commission + fractional\nshares → start with US$5",
      "Mobile-first, simple UX",
      "Financial education content"], PUR),
    ("02", "Long-term\nRetail Investors",
     ["More experienced portfolio builders",
      "USD exposure vs. local currency risk",
      "SIPC protection up to US$500K",
      "ETFs + DRIP for passive wealth"], TEAL),
    ("03", "Crypto\nEnthusiasts",
     ["Tech-savvy, digital-native",
      "Invest from US$1 in crypto",
      "One app: stocks + crypto",
      "Note: crypto not SIPC-covered"], ORANGE),
]

sw11 = 4.1
sg11 = 0.15
for i, (num, title, bullets, c) in enumerate(segments):
    l = 0.35 + i * (sw11 + sg11)
    # Colored header band
    rect(s, l, 2.25, sw11, 0.65, fill=c)
    txt(s, num, l+0.15, 2.3, 0.5, 0.55, size=20, bold=True,
        color=WHITE if c != YELLOW else NEAR_BK)
    txt(s, title, l+0.65, 2.32, sw11-0.8, 0.52,
        size=15, bold=True, color=WHITE if c != YELLOW else NEAR_BK)
    # White body
    rect(s, l, 2.9, sw11, 4.25, fill=WHITE, line=c, lw=1.2)
    for j, b in enumerate(bullets):
        rect(s, l+0.18, 3.1 + j*0.92, 0.38, 0.38, fill=c)
        txt(s, str(j+1), l+0.18, 3.1+j*0.92+0.05, 0.38, 0.3,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, b, l+0.65, 3.1+j*0.92, sw11-0.8, 0.8,
            size=11, color=NEAR_BK)

txt(s, "\"Invest in the world's biggest market from your phone — from US$5, in minutes, with US regulatory protection.\"",
    0.35, 7.15, 12.6, 0.28, size=10, italic=True, color=PUR, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — KEY METRICS & FUNDING
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_slide(s, "Key Metrics & Funding", "06 · Key Metrics & Funding")

# Growth chart (left)
add_image(s, f"{BASE}/chart_growth.png", 0.35, 2.22, 6.5, 4.05)

# Right: stats + funding table
rx12 = 7.1
rect(s, rx12, 2.22, 5.88, 4.05, fill=LIGHT, line=rgb(220,222,235), lw=0.8)
txt(s, "Key Numbers", rx12+0.2, 2.35, 5.4, 0.38, size=13, bold=True, color=NEAR_BK)

kv = [
    ("AUM", "~US$40M", "as of Nov 2023"),
    ("Total Raised", "US$4.3M", "across all rounds"),
    ("Avg. Position", "~US$250", "per user"),
    ("Growth", "10K → 300K", "in 13 months (2022-23)"),
]
for i, (lbl, val, sub) in enumerate(kv):
    t12 = 2.85 + i * 0.86
    rect(s, rx12+0.2, t12, 5.45, 0.78, fill=WHITE, line=rgb(220,222,235), lw=0.5)
    txt(s, lbl + ":", rx12+0.35, t12+0.12, 1.6, 0.3, size=10, color=TEXT_GY)
    txt(s, val, rx12+2.0, t12+0.08, 2.0, 0.38, size=16, bold=True, color=PUR)
    txt(s, sub, rx12+2.0, t12+0.46, 3.0, 0.25, size=9, color=TEXT_GY, italic=True)

# Funding line
rect(s, 0.35, 6.42, 12.6, 0.04, fill=PUR)
txt(s, "Funding", 0.35, 6.55, 2, 0.35, size=12, bold=True, color=NEAR_BK)
fund = [
    ("Sep 2023  ·  US$1.6M", "Utec Ventures · Unpopular Ventures · Softeq Ventures · Mural Capital"),
    ("Total Raised  ·  US$4.3M", "Seed + angel rounds"),
]
for i, (date, inv) in enumerate(fund):
    fl = 0.35 + i * 6.3
    txt(s, date, fl+2.2, 6.55, 3.0, 0.32, size=11, bold=True, color=PUR)
    txt(s, inv, fl+2.2, 6.87, 5.8, 0.28, size=9, color=TEXT_GY, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — SWOT ANALYSIS
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_slide(s, "SWOT Analysis", "07 · SWOT Analysis")

swot = [
    ("STRENGTHS", [
        "Zero commissions + fractional shares from US$5",
        "SEC/FINRA registered + SIPC up to US$500K",
        "20+ LATAM countries · Spanish/Portuguese UX",
        "< 5 min sign-up · mobile-first design",
        "12,000+ instruments + crypto in one app",
    ], GREEN, 0.35, 2.22),
    ("WEAKNESSES", [
        "Low brand awareness outside LATAM",
        "Limited advanced tools for expert traders",
        "Small team (37) — scalability pressure",
        "PFOF model faces regulatory scrutiny",
        "Crypto holdings not covered by SIPC",
    ], RED, 6.8, 2.22),
    ("OPPORTUNITIES", [
        "LATAM fintech market still growing fast",
        "Currency devaluation drives USD-asset demand",
        "Rising smartphone & internet penetration",
        "Huge unbanked/underinvested user base",
        "Expand financial education & content",
    ], TEAL, 0.35, 4.85),
    ("THREATS", [
        "Global brokers entering LATAM (Schwab, IBKR)",
        "Regulatory fragmentation across 20+ countries",
        "PFOF restrictions could cut primary revenue",
        "Local fintechs (Trii, Tyba) accelerating",
        "US market downturns erode user confidence",
    ], ORANGE, 6.8, 4.85),
]

qw = 6.15
qh = 2.35

for title, bullets, c, l, t in swot:
    rect(s, l, t, qw, 0.48, fill=c)
    txt(s, title, l+0.15, t+0.1, qw-0.3, 0.32,
        size=13, bold=True, color=NEAR_BK if c==YELLOW else WHITE)
    rect(s, l, t+0.48, qw, qh-0.48, fill=WHITE, line=c, lw=1.2)
    for i, b in enumerate(bullets):
        txt(s, "▸  " + b, l+0.15, t+0.55+i*0.37, qw-0.3, 0.36,
            size=10.5, color=NEAR_BK)


# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — STRATEGIC OUTLOOK
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
white_slide(s, "Strategic Outlook", "08 · Strategic Outlook")

pillars = [
    ("Scale to\n1M Users",     "Education + referrals\n+ new country launches",  PUR),
    ("Diversify\nRevenue",     "Grow Prime subscriptions\n+ reduce PFOF reliance", TEAL),
    ("Regulatory\nTrust",      "Maintain SEC/FINRA standing\nNavigate multi-country rules", GREEN),
    ("Product\nDepth",         "Add tools for intermediate\n& expert traders", ORANGE),
]

pw = 3.0
pg = 0.16
pl = 0.35
pt14 = 2.3
ph = 2.1

for i, (title, body, c) in enumerate(pillars):
    l = pl + i*(pw+pg)
    rect(s, l, pt14, pw, 0.52, fill=c)
    txt(s, title, l+0.1, pt14+0.06, pw-0.2, 0.44,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, l, pt14+0.52, pw, ph-0.52, fill=WHITE, line=c, lw=1.2)
    txt(s, body, l+0.1, pt14+0.65, pw-0.2, ph-0.75,
        size=12, color=NEAR_BK, align=PP_ALIGN.CENTER)

# Challenge list
rect(s, 0.35, 4.63, 12.6, 0.04, fill=MID_GY)
txt(s, "Challenges to Navigate", 0.35, 4.76, 6, 0.38,
    size=13, bold=True, color=NEAR_BK)

challenges = [
    "Differentiate before global brokers replicate the LATAM-first approach",
    "Continue investing in financial education to convert and retain users",
    "Diversify revenue to reduce exposure to PFOF regulatory risk",
    "Manage multi-country compliance without blowing up cost structure",
]
for i, c14 in enumerate(challenges):
    col = i % 2
    row = i // 2
    cl14 = 0.35 + col * 6.4
    ct14 = 5.22 + row * 0.68
    rect(s, cl14, ct14, 6.2, 0.58, fill=LIGHT, line=rgb(220,222,235), lw=0.5)
    txt(s, "→  " + c14, cl14+0.12, ct14+0.1, 5.95, 0.46, size=11, color=NEAR_BK)

# Quote
txt(s, "\"The platform that wins LATAM earns trust first — through regulation, education, and a UX that feels local.\"",
    0.35, 7.15, 12.6, 0.28, size=10, italic=True, color=PUR, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 15 — THANK YOU
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=DARK)
rect(s, 0, 0, 0.1, 7.5, fill=PUR)
rect(s, 0.1, 0, 0.06, 7.5, fill=TEAL)

# Chart bg at bottom
add_image(s, f"{BASE}/bg_title.png", 0, 4.5, 13.33, 3.0)
rect(s, 0, 4.5, 13.33, 3.0, fill=rgb(30, 35, 60))  # semi-opaque overlay

# Logo
add_image(s, f"{BASE}/hapi_logo.png", 0.55, 0.5, 3.4, 1.02)

rect(s, 0.55, 1.75, 5.5, 0.05, fill=PUR)
txt(s, "Thank You", 0.55, 1.95, 12.2, 1.2,
    size=56, bold=True, color=WHITE)
txt(s, "Questions?", 0.55, 3.1, 12.2, 0.7,
    size=26, color=TEAL, italic=True)

txt(s, "Daniel Taboada  ·  Benjamin Pereyra  ·  Luis Sanchez  ·  Jorge Ramirez  ·  Santiago Wiesse",
    0.55, 5.1, 12.2, 0.45, size=13, color=rgb(180,180,210), align=PP_ALIGN.CENTER)
txt(s, "Universidad del Pacífico  ·  Sem. Fintech  ·  June 2026",
    0.55, 5.6, 12.2, 0.38, size=11, color=rgb(130,130,165),
    italic=True, align=PP_ALIGN.CENTER)
txt(s, "hapi.trade", 0.55, 6.2, 12.2, 0.4,
    size=12, color=TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
OUT = f"{BASE}/Hapi_Case_Study_v2.pptx"
prs.save(OUT)
print(f"\nSaved → {OUT}")
print(f"Slides: {len(prs.slides)}")
