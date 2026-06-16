from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colors
HAPI_PURPLE = RGBColor(0x6B, 0x4E, 0xFF)   # Hapi brand purple
HAPI_TEAL   = RGBColor(0x00, 0xC2, 0xCB)   # accent teal
DARK_BG     = RGBColor(0x0D, 0x1B, 0x2A)   # slide background dark navy
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xF8)
MID_GRAY    = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
GREEN       = RGBColor(0x28, 0xA7, 0x45)
RED         = RGBColor(0xDC, 0x35, 0x45)
YELLOW      = RGBColor(0xFF, 0xC1, 0x07)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # completely blank


# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def bg(slide, color=DARK_BG):
    """Fill entire slide background."""
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=color)


def accent_bar(slide, color=HAPI_PURPLE, height=0.07, top=0.0):
    add_rect(slide, 0, top, 13.33, height, fill_rgb=color)


def section_header(slide, title, subtitle=""):
    bg(slide)
    accent_bar(slide, HAPI_PURPLE)
    accent_bar(slide, HAPI_TEAL, height=0.07, top=7.43)
    add_text(slide, title, 0.5, 2.8, 12.3, 1.2, size=40, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, 0.5, 4.1, 12.3, 0.7, size=18,
                 color=HAPI_TEAL, align=PP_ALIGN.CENTER)


def content_slide_header(slide, title):
    """White bg slide with purple top bar and title."""
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GRAY)
    add_rect(slide, 0, 0, 13.33, 1.0, fill_rgb=DARK_BG)
    add_text(slide, title, 0.4, 0.1, 12.5, 0.8, size=24, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)
    accent_bar(slide, HAPI_PURPLE, height=0.06, top=1.0)


def add_bullet(slide, items, l, t, w, h, size=13, color=DARK_TEXT, spacing=0.38):
    """Draw a list of bullet strings, each on its own line."""
    for i, item in enumerate(items):
        add_text(slide, "▸  " + item, l, t + i * spacing, w, spacing + 0.05,
                 size=size, color=color)


def card(slide, l, t, w, h, fill=WHITE, line=HAPI_PURPLE, lw=1.5):
    r = add_rect(slide, l, t, w, h, fill_rgb=fill, line_rgb=line, line_width=lw)
    return r


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 – TITLE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
# gradient band
add_rect(s, 0, 2.8, 13.33, 2.4, fill_rgb=RGBColor(0x1A, 0x0A, 0x4A))
accent_bar(s, HAPI_PURPLE, height=0.1, top=2.8)
accent_bar(s, HAPI_TEAL,   height=0.1, top=5.1)

add_text(s, "hapi", 0.6, 0.4, 5, 1.2, size=72, bold=True, color=HAPI_PURPLE)
add_text(s, "FINTECH CASE STUDY", 0.6, 1.55, 9, 0.6, size=18, bold=True,
         color=HAPI_TEAL, align=PP_ALIGN.LEFT)

add_text(s, "Democratizing investing in Latin America",
         0.6, 3.0, 12, 0.8, size=28, bold=True, color=WHITE)
add_text(s, "Access to the US market from just US$5",
         0.6, 3.85, 12, 0.5, size=16, color=MID_GRAY, italic=True)

add_text(s, "Daniel Taboada · Benjamin Pereyra · Luis Sanchez\n"
            "Jorge Ramirez · Santiago Wiesse",
         0.6, 5.5, 10, 0.8, size=12, color=MID_GRAY)
add_text(s, "Universidad del Pacífico · Sem. Fintech · June 2026",
         0.6, 6.3, 10, 0.5, size=11, color=RGBColor(0x88, 0x88, 0xAA), italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 – AGENDA
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
content_slide_header(s, "Agenda")

sections = [
    ("01", "Product & Technology"),
    ("02", "Market Opportunity"),
    ("03", "Competitive Landscape"),
    ("04", "Business Model & Revenue"),
    ("05", "Customer Segments"),
    ("06", "Key Metrics & Funding"),
    ("07", "SWOT Analysis"),
    ("08", "Strategic Outlook"),
]

cols = 4
row_h = 1.1
col_w = 3.0
start_l = 0.45
start_t = 1.45

for i, (num, title) in enumerate(sections):
    col = i % cols
    row = i // cols
    l = start_l + col * col_w
    t = start_t + row * (row_h + 0.25)
    card(s, l, t, col_w - 0.15, row_h, fill=WHITE)
    add_text(s, num, l + 0.15, t + 0.05, 0.6, 0.45, size=22, bold=True, color=HAPI_PURPLE)
    add_text(s, title, l + 0.15, t + 0.5, col_w - 0.4, 0.55, size=13,
             bold=False, color=DARK_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 – COMPANY SNAPSHOT
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
content_slide_header(s, "01 · Product & Technology — Company Snapshot")

stats = [
    ("2020", "Founded"),
    ("3", "Co-founders\n(ex UP students)"),
    ("37", "Employees"),
    ("20+", "Countries\nin LATAM"),
    ("500K+", "Users"),
    ("US$5", "Min. investment"),
]

sw = 1.9
gap = 0.12
start = 0.35
top_stat = 1.3

for i, (val, lbl) in enumerate(stats):
    l = start + i * (sw + gap)
    card(s, l, top_stat, sw, 1.5, fill=WHITE)
    add_text(s, val, l + 0.05, top_stat + 0.08, sw - 0.1, 0.75,
             size=28, bold=True, color=HAPI_PURPLE, align=PP_ALIGN.CENTER)
    add_text(s, lbl, l + 0.05, top_stat + 0.82, sw - 0.1, 0.6,
             size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# Story
add_rect(s, 0.35, 3.1, 12.6, 0.05, fill_rgb=HAPI_TEAL)
add_text(s, "The Origin Story", 0.35, 3.25, 5, 0.4, size=14, bold=True, color=DARK_TEXT)
story = ("Three ex-UP students tried to invest in US stocks — found high minimums, "
         "big fees & endless paperwork. So they built their own broker.")
add_text(s, story, 0.35, 3.7, 12.6, 0.6, size=13, color=DARK_TEXT)

add_text(s, "Core Features", 0.35, 4.5, 5, 0.4, size=14, bold=True, color=DARK_TEXT)
features = [
    "Stocks & ETFs (12,000+ instruments) — zero commissions",
    "Fractional shares from US$5  |  Crypto from US$1",
    "DRIP — automatic dividend reinvestment",
    "Hapi Prime subscription (US$9.99/mo): real-time prices, instant settlement, price alerts",
]
add_bullet(s, features, 0.35, 4.95, 12.5, 0.38, size=12)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 – TECHNOLOGY & COMPLIANCE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
content_slide_header(s, "01 · Product & Technology — Infrastructure & Compliance")

# Two columns
# Left
card(s, 0.35, 1.2, 5.9, 5.6, fill=WHITE)
add_text(s, "Technology Stack", 0.5, 1.3, 5.6, 0.5, size=15, bold=True, color=HAPI_PURPLE)
tech_items = [
    "Clearing & Custody → Apex Clearing Corp\n  (same infra as major US brokers)",
    "Crypto → Bakkt Crypto Solutions",
    "2-Factor Authentication + Encryption",
    "Only 37 employees — lean tech team\n  enabled by outsourcing non-core infra",
]
for i, item in enumerate(tech_items):
    add_text(s, "▸  " + item, 0.5, 1.9 + i * 1.1, 5.6, 1.0, size=12, color=DARK_TEXT)

# Right
card(s, 6.6, 1.2, 6.4, 5.6, fill=WHITE)
add_text(s, "Regulatory & Security", 6.75, 1.3, 6.1, 0.5, size=15, bold=True, color=HAPI_PURPLE)

reg_items = [
    ("SEC Registered", "Hapi Securities LLC"),
    ("FINRA Member", "US broker-dealer rules apply"),
    ("SIPC Coverage", "Up to US$500,000 per account"),
    ("Note on Crypto", "Crypto NOT covered by SIPC"),
]
for i, (badge, desc) in enumerate(reg_items):
    t_offset = 1.95 + i * 1.1
    add_rect(s, 6.75, t_offset, 2.2, 0.4, fill_rgb=HAPI_PURPLE)
    add_text(s, badge, 6.8, t_offset + 0.05, 2.1, 0.35, size=11, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, 9.1, t_offset + 0.05, 3.7, 0.35, size=11, color=DARK_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 – MARKET OPPORTUNITY
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
content_slide_header(s, "02 · Market Opportunity")

# Big number
card(s, 0.35, 1.2, 4.0, 2.2, fill=DARK_BG)
add_text(s, "US$15.23B", 0.45, 1.3, 3.8, 1.0, size=30, bold=True,
         color=HAPI_TEAL, align=PP_ALIGN.CENTER)
add_text(s, "LATAM Fintech Market\n(IMARC Group)", 0.45, 2.3, 3.8, 0.8,
         size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

card(s, 4.6, 1.2, 4.0, 2.2, fill=DARK_BG)
add_text(s, "~90%", 4.7, 1.3, 3.8, 1.0, size=30, bold=True,
         color=RED, align=PP_ALIGN.CENTER)
add_text(s, "Banking assets controlled\nby just 5 institutions", 4.7, 2.3,
         3.8, 0.8, size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

card(s, 8.85, 1.2, 4.1, 2.2, fill=DARK_BG)
add_text(s, "1M", 8.95, 1.3, 3.9, 1.0, size=30, bold=True,
         color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, "Users Hapi aims to reach\n(CEO target, Gestión 2025)", 8.95, 2.3,
         3.9, 0.8, size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Trends
add_text(s, "Why now? Key Tailwinds", 0.35, 3.65, 7, 0.4, size=14, bold=True, color=DARK_TEXT)
trends = [
    "Rising smartphone & internet penetration across LATAM",
    "Local currencies losing value → demand for USD-denominated assets",
    "Growing financial literacy and investment culture",
    "Traditional banks ignoring small retail investors",
]
add_bullet(s, trends, 0.35, 4.12, 6.5, 0.38, size=12)

add_text(s, "Hapi's Answer", 7.2, 3.65, 5.8, 0.4, size=14, bold=True, color=HAPI_PURPLE)
answers = [
    "Access US market from US$5",
    "Zero commissions on stocks & ETFs",
    "Sign up in < 5 minutes",
    "Spanish & Portuguese UX + fin-ed content",
]
add_bullet(s, answers, 7.2, 4.12, 5.8, 0.38, size=12)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 – COMPETITIVE LANDSCAPE (table)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
content_slide_header(s, "03 · Competitive Landscape")

headers = ["Platform", "Min. Invest", "Commissions", "SIPC", "LATAM Focus", "Fractional Shares"]
rows = [
    ["Hapi",                "US$5",     "Zero",       "✔ $500K",  "✔ Primary",    "✔"],
    ["Interactive Brokers", "US$0",     "Low",        "✔ $500K",  "Partial",      "✔"],
    ["Charles Schwab",      "US$0",     "Zero",       "✔ $500K",  "✘",            "✔"],
    ["XTB",                 "US$0",     "Zero*",      "✘",        "Partial",      "✔"],
    ["eToro",               "US$10",    "Zero*",      "✘",        "Partial",      "✔"],
    ["Trii / Tyba",         "Low",      "Low",        "✘",        "✔ Local only", "Limited"],
]

col_widths = [2.2, 1.5, 1.7, 1.5, 1.8, 2.0]
col_starts = [0.3]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w + 0.05)

row_h = 0.52
header_t = 1.15
data_t = header_t + row_h + 0.04

# Draw header
for ci, (hdr, cw, cl) in enumerate(zip(headers, col_widths, col_starts)):
    add_rect(s, cl, header_t, cw, row_h, fill_rgb=DARK_BG)
    add_text(s, hdr, cl + 0.05, header_t + 0.08, cw - 0.1, row_h - 0.1,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Draw rows
for ri, row in enumerate(rows):
    row_top = data_t + ri * (row_h + 0.04)
    fill = RGBColor(0xE8, 0xF0, 0xFE) if row[0] == "Hapi" else WHITE
    for ci, (cell, cw, cl) in enumerate(zip(row, col_widths, col_starts)):
        lc = HAPI_PURPLE if row[0] == "Hapi" else MID_GRAY
        lw = 1.5 if row[0] == "Hapi" else 0.5
        add_rect(s, cl, row_top, cw, row_h, fill_rgb=fill, line_rgb=lc, line_width=lw)
        txt_color = HAPI_PURPLE if (row[0] == "Hapi" and ci == 0) else DARK_TEXT
        add_text(s, cell, cl + 0.05, row_top + 0.1, cw - 0.1, row_h - 0.1,
                 size=11, bold=(ci == 0 and row[0] == "Hapi"), color=txt_color,
                 align=PP_ALIGN.CENTER)

add_text(s, "* zero commission on stocks but spreads may apply   ✘ = not applicable for LATAM users",
         0.3, 6.85, 12.7, 0.4, size=9, color=RGBColor(0x77, 0x77, 0x77), italic=True)

add_text(s, "Hapi's Edge: Only platform combining SEC/FINRA/SIPC + zero commissions + LATAM-first UX",
         0.3, 6.45, 12.7, 0.38, size=12, bold=True, color=HAPI_PURPLE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 – BUSINESS MODEL
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
content_slide_header(s, "04 · Business Model & Revenue Streams")

streams = [
    ("Payment for\nOrder Flow", "PFOF — fees from routing\nclient orders to exchanges", HAPI_PURPLE),
    ("Hapi Prime\nSubscription", "US$9.99/mo — real-time\nprices + instant settlement", HAPI_TEAL),
    ("Deposit &\nWithdrawal Fees", "0.9–1% on deposits\nUS$4.99–$10 withdrawals", RGBColor(0xF4, 0x5D, 0x22)),
    ("Interest on\nIdle Cash", "Earns yield on uninvested\ncash in client accounts", GREEN),
    ("Crypto\nSpread", "1% spread on crypto\ntrades via Bakkt", YELLOW),
    ("Closing\nFees", "US$0.10 (whole shares)\nUS$0.15 (fractional)", RGBColor(0x6F, 0x42, 0xC1)),
]

box_w = 2.05
box_h = 1.65
gap_x = 0.07
start_l = 0.25
start_t = 1.3

for i, (title, desc, color) in enumerate(streams):
    l = start_l + i * (box_w + gap_x)
    add_rect(s, l, start_t, box_w, 0.45, fill_rgb=color)
    add_text(s, title, l + 0.05, start_t + 0.04, box_w - 0.1, 0.38,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card(s, l, start_t + 0.45, box_w, box_h - 0.45, fill=WHITE,
         line=color, lw=1.5)
    add_text(s, desc, l + 0.08, start_t + 0.55, box_w - 0.16, box_h - 0.6,
             size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# Key model note
add_rect(s, 0.25, 3.2, 12.85, 0.06, fill_rgb=HAPI_TEAL)
add_text(s, "Business Model Philosophy", 0.25, 3.35, 6, 0.4, size=14, bold=True, color=DARK_TEXT)
notes = [
    "Low-ticket model: avg. position ~US$250 — monetizes through VOLUME & SCALE",
    "Zero-commission label builds trust & acquisition; revenues hidden in PFOF + fees",
    "PFOF controversial: SEC tightened disclosure (Rules 605 & 606); restricted in Europe",
    "Diversified streams reduce dependency on any single revenue source",
]
add_bullet(s, notes, 0.25, 3.82, 12.7, 0.38, size=12)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 – CUSTOMER SEGMENTS
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
content_slide_header(s, "05 · Customer Segments & Value Proposition")

segments = [
    ("Beginner\nInvestors",
     ["Students & young professionals",
      "Want simple, mobile-first UX",
      "Start from US$5 — fractional shares",
      "Zero commissions lower fear of fees"],
     HAPI_PURPLE),
    ("Long-term\nRetail Investors",
     ["More experienced, portfolio builders",
      "Seek USD exposure vs. local currency risk",
      "Value SIPC protection up to US$500K",
      "ETFs + DRIP for passive wealth building"],
     HAPI_TEAL),
    ("Crypto\nEnthusiasts",
     ["Tech-savvy, digital-native users",
      "Invest from US$1 in crypto",
      "Want one app for stocks + crypto",
      "Note: crypto NOT SIPC-covered"],
     RGBColor(0xF4, 0x5D, 0x22)),
]

seg_w = 4.0
seg_h = 5.2
gap = 0.22
start = 0.32
top = 1.25

for i, (title, bullets, color) in enumerate(segments):
    l = start + i * (seg_w + gap)
    add_rect(s, l, top, seg_w, 0.65, fill_rgb=color)
    add_text(s, title, l + 0.1, top + 0.05, seg_w - 0.2, 0.6,
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card(s, l, top + 0.65, seg_w, seg_h - 0.65, fill=WHITE, line=color, lw=1.5)
    for j, b in enumerate(bullets):
        add_text(s, "▸  " + b, l + 0.15, top + 0.82 + j * 0.85,
                 seg_w - 0.3, 0.75, size=12, color=DARK_TEXT)

# Value prop tagline
add_text(s, "\"Invest in the world's biggest market from your phone — from just US$5, in minutes, with US regulatory protection.\"",
         0.32, 6.55, 12.7, 0.5, size=13, italic=True, color=HAPI_PURPLE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 – KEY METRICS & FUNDING
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
content_slide_header(s, "06 · Key Metrics, Unit Economics & Funding")

# Timeline of growth
add_text(s, "User Growth Timeline", 0.35, 1.2, 8, 0.4, size=14, bold=True, color=DARK_TEXT)

timeline = [
    ("Oct 2022", "10,000\nusers"),
    ("Nov 2023", "300,000\nusers"),
    ("2025", "500,000+\nusers"),
    ("Goal", "1,000,000\nusers"),
]
tw = 2.8
tl_top = 1.75
for i, (date, val) in enumerate(timeline):
    l = 0.35 + i * (tw + 0.1)
    clr = HAPI_PURPLE if i < 3 else HAPI_TEAL
    add_rect(s, l, tl_top, tw, 0.38, fill_rgb=clr)
    add_text(s, date, l + 0.05, tl_top + 0.05, tw - 0.1, 0.3,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, l, tl_top + 0.38, tw, 0.85, fill_rgb=WHITE,
             line_rgb=clr, line_width=1.5)
    add_text(s, val, l + 0.05, tl_top + 0.43, tw - 0.1, 0.7,
             size=18, bold=True, color=clr, align=PP_ALIGN.CENTER)

# Right stats
stats2 = [
    ("AUM", "~US$40M", "(as of Nov 2023)"),
    ("Avg. Position", "~US$250", "per user"),
    ("Employees", "37", "lean team"),
]
for i, (lbl, val, sub) in enumerate(stats2):
    l = 9.5
    t = 1.2 + i * 1.4
    card(s, l, t, 3.5, 1.2, fill=DARK_BG)
    add_text(s, lbl, l + 0.1, t + 0.06, 3.2, 0.35, size=11, color=MID_GRAY)
    add_text(s, val, l + 0.1, t + 0.38, 3.2, 0.5, size=22, bold=True, color=HAPI_TEAL)
    add_text(s, sub, l + 0.1, t + 0.88, 3.2, 0.25, size=10, color=MID_GRAY, italic=True)

# Funding
add_rect(s, 0.35, 3.35, 8.8, 0.06, fill_rgb=HAPI_PURPLE)
add_text(s, "Funding Rounds", 0.35, 3.5, 5, 0.4, size=14, bold=True, color=DARK_TEXT)

fund_headers = ["Date", "Amount", "Investors"]
fund_rows = [
    ["Pre-seed / Seed", "US$2.7M",  "Various angel / early investors"],
    ["Sep 2023",        "US$1.6M",  "Utec Ventures, Unpopular Ventures,\nSofteq Ventures, Mural Capital"],
    ["Total Raised",    "US$4.3M",  "—"],
]
fh = [1.8, 1.3, 5.5]
fl = [0.35, 2.25, 3.65]
ft = 4.0

for ci, (h, w, l) in enumerate(zip(fund_headers, fh, fl)):
    add_rect(s, l, ft, w, 0.4, fill_rgb=DARK_BG)
    add_text(s, h, l + 0.05, ft + 0.08, w - 0.1, 0.28,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(fund_rows):
    rt = ft + 0.44 + ri * 0.72
    fill = RGBColor(0xE8, 0xF0, 0xFE) if row[0] == "Total Raised" else WHITE
    for ci, (cell, w, l) in enumerate(zip(row, fh, fl)):
        add_rect(s, l, rt, w, 0.65, fill_rgb=fill,
                 line_rgb=MID_GRAY, line_width=0.5)
        bold = (row[0] == "Total Raised")
        add_text(s, cell, l + 0.05, rt + 0.08, w - 0.1, 0.55,
                 size=11, bold=bold, color=DARK_TEXT if not bold else HAPI_PURPLE,
                 align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 – SWOT
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
content_slide_header(s, "07 · SWOT Analysis")

quadrants = [
    ("STRENGTHS", [
        "Zero-commission + fractional shares from US$5",
        "SEC/FINRA registered + SIPC coverage up to US$500K",
        "20+ LATAM countries — Spanish & Portuguese UX",
        "Simple onboarding: < 5 min sign-up",
        "12,000+ instruments including ETFs & crypto",
    ], GREEN, 0.35, 1.15),
    ("WEAKNESSES", [
        "Low brand awareness outside LATAM",
        "Lacks advanced tools for expert traders",
        "Small team (37) — scalability pressure",
        "PFOF model faces regulatory scrutiny",
        "Crypto not covered by SIPC",
    ], RED, 6.85, 1.15),
    ("OPPORTUNITIES", [
        "LATAM fintech market still growing",
        "Currency devaluation drives USD demand",
        "Rising smartphone & internet penetration",
        "Low investment penetration → huge user base",
        "Expand financial education content",
    ], HAPI_TEAL, 0.35, 4.25),
    ("THREATS", [
        "Global brokers (Schwab, IBKR) entering LATAM",
        "Regulatory fragmentation across 20+ countries",
        "PFOF restrictions could cut main revenue",
        "Local fintechs (Trii, Tyba) competing",
        "Macro risk: US market downturns hit user trust",
    ], YELLOW, 6.85, 4.25),
]

qw = 6.2
qh = 2.85
for title, bullets, color, l, t in quadrants:
    add_rect(s, l, t, qw, 0.45, fill_rgb=color)
    add_text(s, title, l + 0.1, t + 0.06, qw - 0.2, 0.35,
             size=13, bold=True, color=DARK_TEXT if color == YELLOW else WHITE,
             align=PP_ALIGN.LEFT)
    card(s, l, t + 0.45, qw, qh - 0.45, fill=WHITE, line=color, lw=1.2)
    for i, b in enumerate(bullets):
        add_text(s, "▸  " + b, l + 0.15, t + 0.55 + i * 0.44,
                 qw - 0.3, 0.42, size=11, color=DARK_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 – STRATEGIC OUTLOOK
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
content_slide_header(s, "08 · Strategic Outlook")

pillars = [
    ("Scale Users",       "Reach 1M users\nvia education + referrals",  HAPI_PURPLE),
    ("Deepen Revenue",    "Grow Prime subscriptions\n& reduce PFOF dependency", HAPI_TEAL),
    ("Regulatory Trust",  "Maintain SEC/FINRA standing\nNavigate multi-country rules", GREEN),
    ("Product Depth",     "Add advanced tools\nfor intermediate traders", RGBColor(0xF4, 0x5D, 0x22)),
]

pw = 2.95
ph = 2.0
pl_start = 0.35
pt = 1.3
gap_p = 0.17

for i, (title, body, color) in enumerate(pillars):
    l = pl_start + i * (pw + gap_p)
    add_rect(s, l, pt, pw, 0.45, fill_rgb=color)
    add_text(s, title, l + 0.1, pt + 0.06, pw - 0.2, 0.36,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card(s, l, pt + 0.45, pw, ph - 0.45, fill=WHITE, line=color, lw=1.5)
    add_text(s, body, l + 0.1, pt + 0.55, pw - 0.2, ph - 0.6,
             size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# Key challenges
add_rect(s, 0.35, 3.55, 12.65, 0.06, fill_rgb=MID_GRAY)
add_text(s, "Key Challenges to Navigate", 0.35, 3.7, 8, 0.4, size=14, bold=True, color=DARK_TEXT)
challenges = [
    "Differentiate before global brokers replicate the LATAM-first approach",
    "Continue investing in financial education to retain & convert users",
    "Diversify revenue to reduce exposure to PFOF regulatory risk",
    "Manage multi-country compliance without blowing up the cost structure",
]
add_bullet(s, challenges, 0.35, 4.18, 12.65, 0.4, size=12)

add_text(s,
    "\"The platform that wins LATAM will be the one that earns trust first — through regulation, education, and a UX that feels local.\"",
    0.35, 6.45, 12.65, 0.6, size=13, italic=True, color=HAPI_PURPLE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 – THANK YOU
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
accent_bar(s, HAPI_PURPLE, height=0.12, top=0.0)
accent_bar(s, HAPI_TEAL,   height=0.12, top=7.38)

add_text(s, "hapi", 0.5, 1.8, 12.3, 1.5, size=80, bold=True,
         color=HAPI_PURPLE, align=PP_ALIGN.CENTER)
add_text(s, "Thank You", 0.5, 3.4, 12.3, 0.8, size=32, bold=False,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Questions?", 0.5, 4.1, 12.3, 0.6, size=22, bold=False,
         color=HAPI_TEAL, align=PP_ALIGN.CENTER)
add_text(s, "Daniel Taboada · Benjamin Pereyra · Luis Sanchez · Jorge Ramirez · Santiago Wiesse",
         0.5, 5.3, 12.3, 0.5, size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Universidad del Pacífico — Sem. Fintech — June 2026",
         0.5, 5.85, 12.3, 0.4, size=11, color=RGBColor(0x77, 0x77, 0x99),
         italic=True, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
OUT = "/home/user/clase13/Hapi_Case_Study.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
